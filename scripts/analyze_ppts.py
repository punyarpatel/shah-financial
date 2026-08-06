import sys
import os
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

def analyze_ppt(file_path):
    if not os.path.exists(file_path):
        return f"File not found: {file_path}"
    
    prs = Presentation(file_path)
    slide_count = len(prs.slides)
    slides_summary = []

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        texts = []
        has_images = False
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        texts.append(t)
            if shape.shape_type == 13: # Picture shape
                has_images = True

        title = texts[0] if texts else "Untitled Slide"
        slides_summary.append({
            "slide_num": slide_num,
            "title": title,
            "text_count": len(texts),
            "sample_texts": texts[:6],
            "has_images": has_images
        })

    return {
        "file_name": os.path.basename(file_path),
        "total_slides": slide_count,
        "slides": slides_summary
    }

def run_analysis():
    p1 = r"D:\vite-project\public\Drishti_Wealth_Internship_Deck_v2.pptx"
    p2 = r"D:\vite-project\public\PPT Gemini 2.pptx"

    print("=== PPT 1: Drishti_Wealth_Internship_Deck_v2.pptx ===")
    res1 = analyze_ppt(p1)
    print(f"Total Slides: {res1['total_slides']}")
    for s in res1['slides']:
        print(f"Slide {s['slide_num']}: {s['title']} | Has Images: {s['has_images']}")
        for t in s['sample_texts'][1:4]:
            print(f"   - {t[:80]}")

    print("\n=== PPT 2: PPT Gemini 2.pptx ===")
    res2 = analyze_ppt(p2)
    print(f"Total Slides: {res2['total_slides']}")
    for s in res2['slides']:
        print(f"Slide {s['slide_num']}: {s['title']} | Has Images: {s['has_images']}")
        for t in s['sample_texts'][1:4]:
            print(f"   - {t[:80]}")

if __name__ == "__main__":
    run_analysis()
